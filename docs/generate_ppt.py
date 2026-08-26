"""작업자 안내용 PPT 생성 (간결 6장)"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SCRIPT_DIR = Path(__file__).resolve().parent
OUT = SCRIPT_DIR / "휴가증 자동 반영 프로그램 안내 (작업자용).pptx"

# 본문 폭과 왼쪽 끝. 슬라이드 가운데(6.6665")를 기준으로 잡는다.
#
# 예전에는 왼쪽 0.8 에서 시작해 간격만 곱했는데, 마지막 카드가 13.10 에서 끝나
# 오른쪽 여백이 0.23 밖에 안 남았다 — 눈에 띄게 오른쪽으로 밀려 보인다.
# 폭을 먼저 정하고 남는 만큼을 좌우로 나누면 카드 개수와 무관하게 가운데다.
SLIDE_W = 13.333
CONTENT_W = 11.9
CONTENT_X = (SLIDE_W - CONTENT_W) / 2      # 0.7165


def cards(n, gap=0.3):
    """카드 n개를 본문 폭 안에 고르게 편다. (왼쪽 x 목록, 카드 폭) 반환."""
    w = (CONTENT_W - gap * (n - 1)) / n
    return [CONTENT_X + i * (w + gap) for i in range(n)], w

# COSMAX 색상
RED = RGBColor(0xC8, 0x10, 0x2E)
DARK_RED = RGBColor(0xA3, 0x0D, 0x24)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x77, 0x77, 0x77)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_RED = RGBColor(0xFF, 0xF5, 0xF5)


def set_text(tf, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "맑은 고딕"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_rounded_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.12
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = anchor
    set_text(tf, text, size=size, bold=bold, color=color, align=align)
    return tf


def page_footer(slide, n, total):
    add_text(slide, Inches(0.4), Inches(7.0), Inches(12.6), Inches(0.3),
             f"COSMAX · 생산3팀 파우더 성형실    |    {n} / {total}",
             size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    TOTAL = 6

    # ===== 1. 표지 =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, Inches(0.4), prs.slide_height, RED)
    add_text(s, Inches(1.0), Inches(1.5), Inches(11.5), Inches(0.5),
             "COSMAX · 생산3팀 파우더 성형실", size=16, bold=True, color=RED)
    add_text(s, Inches(1.0), Inches(2.4), Inches(11.5), Inches(1.5),
             "휴가증 자동 반영 프로그램", size=48, bold=True, color=DARK)
    add_text(s, Inches(1.0), Inches(4.0), Inches(11.5), Inches(0.6),
             "현장 작업자 안내", size=22, color=GRAY)
    add_rect(s, Inches(1.0), Inches(4.7), Inches(2.0), Pt(3), RED)
    add_text(s, Inches(1.0), Inches(5.5), Inches(11.5), Inches(1.0),
             "본인 휴대폰으로 직접 휴가증을 작성하세요.", size=18, color=DARK)
    add_text(s, Inches(1.0), Inches(6.7), Inches(11.5), Inches(0.4),
             "발행: 2026-05-28  ·  개정: 2026-08-26  ·  담당: 이동준", size=11, color=GRAY)

    # ===== 2. 어떻게 바뀌나요 =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, prs.slide_width, Inches(1.0), RED)
    add_text(s, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6),
             "어떻게 바뀌나요?", size=28, bold=True, color=WHITE)
    xs2, cw2 = cards(2)
    # 이전
    add_rounded_rect(s, Inches(xs2[0]), Inches(1.7), Inches(cw2), Inches(4.3), LIGHT_GRAY)
    add_text(s, Inches(xs2[0] + 0.2), Inches(2.0), Inches(cw2 - 0.4), Inches(0.6),
             "❌  이전 방식", size=18, bold=True, color=GRAY)
    add_text(s, Inches(xs2[0] + 0.2), Inches(2.9), Inches(cw2 - 0.4), Inches(3.0),
             "• 종이 휴가증 작성\n\n• 휴가증함에 제출\n\n• 담당자가 매일 수거\n\n• 서무가 일일이 입력",
             size=15, color=DARK)
    # 새 방식
    add_rounded_rect(s, Inches(xs2[1]), Inches(1.7), Inches(cw2), Inches(4.3), LIGHT_RED)
    add_text(s, Inches(xs2[1] + 0.2), Inches(2.0), Inches(cw2 - 0.4), Inches(0.6),
             "✓  새 방식", size=18, bold=True, color=RED)
    add_text(s, Inches(xs2[1] + 0.2), Inches(2.9), Inches(cw2 - 0.4), Inches(3.0),
             "• 본인 휴대폰으로 사이트 접속\n\n• 본인 휴가증 직접 작성\n\n• 클라우드에 자동 저장\n\n• 서무가 한 번에 자동 등록",
             size=15, color=DARK)
    add_text(s, Inches(CONTENT_X), Inches(6.3), Inches(CONTENT_W), Inches(0.5),
             "→ 작성자는 그대로 본인. 종이 대신 휴대폰으로!",
             size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)
    page_footer(s, 2, TOTAL)

    # ===== 3. 사이트 접속 & 로그인 =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, prs.slide_width, Inches(1.0), RED)
    add_text(s, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6),
             "1. 접속 & 로그인", size=28, bold=True, color=WHITE)

    # 사이트 주소
    add_rounded_rect(s, Inches(CONTENT_X), Inches(1.5), Inches(CONTENT_W), Inches(0.9), LIGHT_RED)
    add_text(s, Inches(CONTENT_X), Inches(1.5), Inches(CONTENT_W), Inches(0.9),
             "https://djlee-cosmax.github.io/Vacation-Permit-Auto-System/",
             size=17, bold=True, color=RED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 단계 3개
    steps = [
        ("①", "휴대폰으로 접속", "Safari / Chrome"),
        ("②", "사번 + 1234", "초기 비밀번호"),
        ("③", "로그인", "→ 메인 화면 진입"),
    ]
    xs, cw = cards(3)
    for (num, title, desc), cx in zip(steps, xs):
        x = Inches(cx)
        add_rounded_rect(s, x, Inches(2.8), Inches(cw), Inches(2.6), LIGHT_RED)
        add_text(s, x, Inches(3.0), Inches(cw), Inches(0.7),
                 num, size=36, bold=True, color=RED, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(3.9), Inches(cw), Inches(0.5),
                 title, size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(4.5), Inches(cw), Inches(0.6),
                 desc, size=13, color=GRAY, align=PP_ALIGN.CENTER)

    # 안내 — 보안 질문은 폐지됐다(2026-07 인증 서버 이관). 대신 자동 로그아웃을 알린다.
    # 푸터가 7.0" 에 있어 상자는 6.95" 안에서 끝나야 한다.
    add_rounded_rect(s, Inches(CONTENT_X), Inches(5.55), Inches(CONTENT_W), Inches(1.35), LIGHT_GRAY)
    add_text(s, Inches(CONTENT_X + 0.3), Inches(5.65), Inches(CONTENT_W - 0.6), Inches(1.15),
             "⚠️ 첫 로그인 후 [내 정보] → [비밀번호 변경] → 새 비밀번호(숫자 6~10자리)\n"
             "⏱ 10분 동안 조작이 없으면 자동 로그아웃됩니다 — 다시 로그인하면 됩니다.",
             size=14, bold=True, color=DARK_RED, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 3, TOTAL)

    # ===== 4. 홈 화면에 추가 (PWA) =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, prs.slide_width, Inches(1.0), RED)
    add_text(s, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6),
             "2. 홈 화면에 추가하기 (권장)", size=28, bold=True, color=WHITE)

    add_text(s, Inches(CONTENT_X), Inches(1.4), Inches(CONTENT_W), Inches(0.6),
             "휴대폰 홈 화면에 추가하면 일반 앱처럼 아이콘 한 번에 실행됩니다.",
             size=15, color=DARK, align=PP_ALIGN.CENTER)

    # iOS 카드
    xs4, cw4 = cards(2)
    add_rounded_rect(s, Inches(xs4[0]), Inches(2.3), Inches(cw4), Inches(3.8), LIGHT_RED)
    add_text(s, Inches(xs4[0] + 0.2), Inches(2.5), Inches(cw4 - 0.4), Inches(0.6),
             "📱 iOS (Safari)", size=20, bold=True, color=RED)
    add_text(s, Inches(xs4[0] + 0.2), Inches(3.3), Inches(cw4 - 0.4), Inches(2.7),
             "① 사이트 접속\n\n"
             "② 하단 가운데 공유 버튼 ⬆️ 탭\n\n"
             "③ \"홈 화면에 추가\" 선택\n\n"
             "④ 우상단 \"추가\" 탭",
             size=14, color=DARK)

    # Android 카드
    add_rounded_rect(s, Inches(xs4[1]), Inches(2.3), Inches(cw4), Inches(3.8), LIGHT_RED)
    add_text(s, Inches(xs4[1] + 0.2), Inches(2.5), Inches(cw4 - 0.4), Inches(0.6),
             "🤖 Android (Chrome)", size=20, bold=True, color=RED)
    add_text(s, Inches(xs4[1] + 0.2), Inches(3.3), Inches(cw4 - 0.4), Inches(2.7),
             "① 사이트 접속\n\n"
             "② 우상단 ⋮ (점 세 개) 탭\n\n"
             "③ \"홈 화면에 추가\" 선택\n\n"
             "④ \"설치\" 또는 \"추가\" 확인",
             size=14, color=DARK)

    # 안내 박스
    add_rounded_rect(s, Inches(CONTENT_X), Inches(6.3), Inches(CONTENT_W), Inches(0.7), LIGHT_GRAY)
    add_text(s, Inches(CONTENT_X + 0.2), Inches(6.3), Inches(CONTENT_W - 0.4), Inches(0.7),
             "💡 추가 후 아이콘 탭 → 풀스크린 앱으로 실행 / 새로고침은 상단 ↻ 아이콘",
             size=12, bold=True, color=DARK_RED, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 4, TOTAL)

    # ===== 5. 휴가증 작성 + 유형 =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, prs.slide_width, Inches(1.0), RED)
    add_text(s, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6),
             "3. 휴가증 작성", size=28, bold=True, color=WHITE)

    # 좌측 - 작성 절차
    add_text(s, Inches(CONTENT_X), Inches(1.5), Inches(6.0), Inches(0.5),
             "왼쪽 [휴가증 작성] 카드에서:", size=14, bold=True, color=DARK)
    steps = [
        ("1.", "구분 선택", "연차, 반차, 생휴 등"),
        ("2.", "개수 + 기간", "예: 연차 2개"),
        ("3.", "사유 입력", "휴가 사유"),
        ("4.", "[+ 추가]", "클라우드 자동 저장"),
    ]
    for i, (n, t, d) in enumerate(steps):
        y = Inches(2.1 + i * 0.7)
        add_text(s, Inches(CONTENT_X), y, Inches(0.4), Inches(0.5),
                 n, size=15, bold=True, color=RED)
        add_text(s, Inches(CONTENT_X + 0.5), y, Inches(2.2), Inches(0.5),
                 t, size=14, bold=True, color=DARK)
        add_text(s, Inches(CONTENT_X + 2.8), y, Inches(3.2), Inches(0.5),
                 d, size=12, color=GRAY)
    # 안내
    add_rounded_rect(s, Inches(CONTENT_X), Inches(5.05), Inches(6.0), Inches(1.75), LIGHT_RED)
    add_text(s, Inches(CONTENT_X + 0.2), Inches(5.18), Inches(5.6), Inches(1.5),
             "• 이름·연락처는 자동 채움 (수정 불가)\n"
             "• 여러 유형은 휴가증 따로 작성 (예: 연차+반차 = 2건)\n"
             "• 같은 날짜 합계 1일 초과 시 등록 차단\n"
             "• 잔여보다 많이 신청하면 작성 차단\n"
             "• 서무가 처리 완료한 휴가증은 재등록 불가\n"
             "• 하기휴가는 1개 = 연속 3일 — 종료일 자동, 한 장으로 작성",
             size=12, color=DARK)

    # 우측 - 휴가 유형 표
    add_text(s, Inches(CONTENT_X + 6.4), Inches(1.5), Inches(5.5), Inches(0.5),
             "휴가 유형:", size=14, bold=True, color=DARK)
    rows = [
        ("구분", "1개당"),
        ("연차", "1.0일"),
        ("반차(오전/오후)", "0.5일"),
        ("반반차(오전/오후)", "0.25일"),
        ("생휴", "1.0일"),
        ("하기휴가", "3.0일"),
        ("결근(전/오전/오후)", "1.0/0.5일"),
    ]
    tbl_x = Inches(CONTENT_X + 6.4); tbl_y = Inches(2.1); rh = Inches(0.55)
    cols = [Inches(3.5), Inches(2.0)]
    for ri, row in enumerate(rows):
        y = tbl_y + rh * ri
        is_header = (ri == 0)
        bg = RED if is_header else (LIGHT_GRAY if ri % 2 == 0 else WHITE)
        add_rect(s, tbl_x, y, Inches(5.5), rh, bg)
        x = tbl_x
        for ci, txt in enumerate(row):
            add_text(s, x + Inches(0.2), y, cols[ci] - Inches(0.2), rh,
                     txt, size=12, bold=is_header,
                     color=WHITE if is_header else DARK, anchor=MSO_ANCHOR.MIDDLE)
            x += cols[ci]
    page_footer(s, 5, TOTAL)

    # ===== 6. 내 휴가증 & 자주 묻는 질문 =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, prs.slide_width, Inches(1.0), RED)
    add_text(s, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6),
             "4. 내 휴가증 확인 & FAQ", size=28, bold=True, color=WHITE)

    # 내 휴가증 (좌측)
    add_text(s, Inches(CONTENT_X), Inches(1.4), Inches(6.0), Inches(0.5),
             "내 휴가증 확인", size=16, bold=True, color=RED)
    add_rounded_rect(s, Inches(CONTENT_X), Inches(2.0), Inches(6.0), Inches(2.4), LIGHT_RED)
    add_text(s, Inches(CONTENT_X + 0.2), Inches(2.15), Inches(5.6), Inches(2.2),
             "① 상단 [내 휴가증] 클릭\n\n"
             "② 서무가 등록 완료한 본인 휴가증이\n   바로 표시됩니다.\n   (✓ 처리 완료 표시, 최근 14일)\n\n"
             "③ 위쪽에 내 잔여 휴가(개수)도 함께 표시\n   연차 · 생휴 · 하기휴가(5~10월)",
             size=13, color=DARK)
    # 세 줄짜리 안내다. 상자를 2.0" 로 두면 아래가 텅 빈 회색 덩어리로 보인다.
    add_rounded_rect(s, Inches(CONTENT_X), Inches(4.75), Inches(6.0), Inches(1.25), LIGHT_GRAY)
    add_text(s, Inches(CONTENT_X + 0.2), Inches(4.85), Inches(5.6), Inches(1.05),
             "⚠️ 잘못 작성 시\n작성 직후 우측 카드의 [삭제]로 즉시 취소.\n서무 처리 후엔 별도 요청 필요.",
             size=12, color=DARK)

    # FAQ (우측)
    add_text(s, Inches(CONTENT_X + 6.4), Inches(1.4), Inches(5.5), Inches(0.5),
             "자주 묻는 질문", size=16, bold=True, color=RED)
    # [비밀번호 찾기] 는 없어졌다 — 보안 질문 폐지(2026-07)로 서무 요청 → 관리자 처리 방식이다.
    qa = [
        ("Q. 비밀번호 분실?", "서무에게 초기화 요청\n     (관리자 처리 후 1234로 로그인)"),
        ("Q. 잔여가 부족하다고 나옴?", "[내 휴가증] 상단에서 남은 개수 확인"),
        ("Q. 사번 등록 안 됨?", "관리자(이동준)에게 요청"),
        ("Q. 화면이 옛 디자인?", "상단 ↻ 아이콘 클릭"),
        ("Q. 서버 저장 실패?", "인터넷 확인 후 재작성"),
    ]
    for i, (q, a) in enumerate(qa):
        y = Inches(1.95 + i * 0.95)
        add_text(s, Inches(CONTENT_X + 6.4), y, Inches(5.5), Inches(0.35),
                 q, size=13, bold=True, color=DARK)
        add_text(s, Inches(CONTENT_X + 6.6), y + Inches(0.35), Inches(5.3), Inches(0.55),
                 "→ " + a, size=12, color=GRAY)
    page_footer(s, 6, TOTAL)

    prs.save(str(OUT))
    print(f"PPT 생성 완료: {OUT}")
    print(f"크기: {OUT.stat().st_size:,} bytes")


if __name__ == "__main__":
    main()
