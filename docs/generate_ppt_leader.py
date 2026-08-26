"""서무 담당자용 안내 PPT 생성 (6장).

작업자용(generate_ppt.py)과 색·글꼴·여백을 맞춘다 — 도우미 함수를 그대로 쓴다.
generate_ppt.py 는 __main__ 가드가 있어 import 해도 파일을 만들지 않는다.

이 자료의 중심은 3장이다. 2026-08-26 부터 잔여 차감을 서버가 하므로 서무가
손으로 맞추지 않아도 된다 — 지금까지 해 오던 일이 하나 없어진 셈이라
그것부터 분명히 해야 한다.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from generate_ppt import (
    RED, DARK_RED, DARK, GRAY, LIGHT_GRAY, WHITE, LIGHT_RED,
    add_rect, add_rounded_rect, add_text,
)

OUT = Path(__file__).resolve().parent / "휴가증 자동 반영 프로그램 안내 (서무용).pptx"
TOTAL = 6


def header(prs, slide, title):
    add_rect(slide, 0, 0, prs.slide_width, Inches(1.0), RED)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6),
             title, size=28, bold=True, color=WHITE)


def footer(slide, n):
    add_text(slide, Inches(0.4), Inches(7.0), Inches(12.6), Inches(0.3),
             f"COSMAX · 생산3팀 파우더 성형실    |    {n} / {TOTAL}",
             size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ===== 1. 표지 =====
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, Inches(0.4), prs.slide_height, RED)
    add_text(s, Inches(1.0), Inches(1.5), Inches(11.5), Inches(0.5),
             "COSMAX · 생산3팀 파우더 성형실", size=16, bold=True, color=RED)
    add_text(s, Inches(1.0), Inches(2.4), Inches(11.5), Inches(1.5),
             "휴가증 자동 반영 프로그램", size=48, bold=True, color=DARK)
    add_text(s, Inches(1.0), Inches(4.0), Inches(11.5), Inches(0.6),
             "서무 담당자 안내", size=22, color=GRAY)
    add_rect(s, Inches(1.0), Inches(4.7), Inches(2.0), Pt(3), RED)
    add_text(s, Inches(1.0), Inches(5.5), Inches(11.5), Inches(1.0),
             "매일 하는 일은 세 번의 클릭입니다.", size=18, color=DARK)
    add_text(s, Inches(1.0), Inches(6.7), Inches(11.5), Inches(0.4),
             "발행: 2026-08-26  ·  담당: 이동준", size=11, color=GRAY)

    # ===== 2. 매일 하는 일 =====
    s = prs.slides.add_slide(blank)
    header(prs, s, "1. 매일 하는 일")

    add_text(s, Inches(0.8), Inches(1.25), Inches(11.7), Inches(0.5),
             "화면 오른쪽 위 버튼을 ① ② ③ 순서대로 누르면 끝입니다.",
             size=15, color=DARK, align=PP_ALIGN.CENTER)

    steps = [
        ("①", "서버에서 불러오기", "작업자가 올린 휴가증을\n가져옵니다"),
        ("②", "파일로 내보내기", "JSON · XLSX 가\n다운로드 폴더에 저장"),
        ("③", "프로그램 실행", "자동화가 그룹웨어에\n임시저장합니다"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.8 + i * 4.2)
        add_rounded_rect(s, x, Inches(1.95), Inches(3.9), Inches(2.9), LIGHT_RED)
        add_text(s, x, Inches(2.15), Inches(3.9), Inches(0.7),
                 num, size=36, bold=True, color=RED, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(3.0), Inches(3.9), Inches(0.5),
                 title, size=17, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(3.6), Inches(3.9), Inches(1.1),
                 desc, size=12.5, color=GRAY, align=PP_ALIGN.CENTER)

    add_rounded_rect(s, Inches(0.8), Inches(5.15), Inches(11.9), Inches(0.85), LIGHT_GRAY)
    add_text(s, Inches(1.1), Inches(5.15), Inches(11.3), Inches(0.85),
             "④ 그룹웨어에서 내용을 검토하고 직접 신청합니다 — 마지막 확인은 사람이 합니다.",
             size=13.5, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

    add_rounded_rect(s, Inches(0.8), Inches(6.15), Inches(11.9), Inches(0.75), LIGHT_RED)
    add_text(s, Inches(1.1), Inches(6.15), Inches(11.3), Inches(0.75),
             "⚠️ run.bat 을 직접 열지 마세요. ③ [프로그램 실행] 을 눌러야 같은 자동화가 돕니다.",
             size=13.5, bold=True, color=DARK_RED, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 2)

    # ===== 3. 잔여 휴가는 손대지 않습니다 (핵심) =====
    s = prs.slides.add_slide(blank)
    header(prs, s, "2. 잔여 휴가는 손대지 않습니다")

    add_text(s, Inches(0.8), Inches(1.25), Inches(11.7), Inches(0.6),
             "2026년 8월부터 잔여 차감을 서버가 매일 저녁 자동으로 합니다.",
             size=17, bold=True, color=RED, align=PP_ALIGN.CENTER)

    # 이전 / 지금
    add_rounded_rect(s, Inches(0.8), Inches(2.1), Inches(5.7), Inches(2.5), LIGHT_GRAY)
    add_text(s, Inches(1.0), Inches(2.3), Inches(5.3), Inches(0.5),
             "지금까지", size=17, bold=True, color=GRAY)
    add_text(s, Inches(1.0), Inches(2.95), Inches(5.3), Inches(1.5),
             "• 잔여가 안 맞으면\n   명단에서 손으로 고침\n• 매번 확인이 필요했음",
             size=14, color=DARK)

    add_rounded_rect(s, Inches(6.8), Inches(2.1), Inches(5.7), Inches(2.5), LIGHT_RED)
    add_text(s, Inches(7.0), Inches(2.3), Inches(5.3), Inches(0.5),
             "앞으로", size=17, bold=True, color=RED)
    add_text(s, Inches(7.0), Inches(2.95), Inches(5.3), Inches(1.5),
             "• 서버가 매일 저녁 자동 차감\n• 서무는 아무것도 안 해도 됨\n• 시간이 지나면 저절로 맞음",
             size=14, color=DARK)

    # 예외
    add_rounded_rect(s, Inches(0.8), Inches(4.85), Inches(11.9), Inches(1.0), LIGHT_GRAY)
    add_text(s, Inches(1.1), Inches(4.85), Inches(11.3), Inches(1.0),
             "다만 연차는 그룹웨어 실제 잔여와 다르면 고쳐도 됩니다.\n"
             "고친 시점이 새 기준이 되어 그 이전 휴가증은 다시 빠지지 않습니다.",
             size=13.5, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

    add_rounded_rect(s, Inches(0.8), Inches(6.0), Inches(11.9), Inches(0.9), LIGHT_RED)
    add_text(s, Inches(1.1), Inches(6.0), Inches(11.3), Inches(0.9),
             "⚠️ 생휴 · 하기휴가 칸은 절대 손대지 마세요.\n"
             "     손으로 줄이면 자동 차감이 또 빼서 마이너스가 됩니다.",
             size=13.5, bold=True, color=DARK_RED, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 3)

    # ===== 4. 작업자 명단 =====
    s = prs.slides.add_slide(blank)
    header(prs, s, "3. 작업자 명단")

    add_text(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5),
             "상단 [작업자 명단] 을 누르면 전체 명단과 잔여 휴가가 보입니다.",
             size=15, color=DARK, align=PP_ALIGN.CENTER)

    can = [
        ("할 수 있는 것", RED, LIGHT_RED,
         "• 잔여 연차 수정 후 [저장]\n"
         "• 이름 · 사번 · 근무지 확인\n"
         "• 비밀번호 [초기화 요청]"),
        ("관리자에게 요청할 것", GRAY, LIGHT_GRAY,
         "• 신규 입사자 명단 추가\n"
         "• 퇴직자 명단 삭제\n"
         "• 비밀번호 실제 초기화 처리"),
    ]
    for i, (title, col, bg, body) in enumerate(can):
        x = Inches(0.8 + i * 6.0)
        add_rounded_rect(s, x, Inches(2.0), Inches(5.7), Inches(2.9), bg)
        add_text(s, x + Inches(0.25), Inches(2.2), Inches(5.2), Inches(0.5),
                 title, size=17, bold=True, color=col)
        add_text(s, x + Inches(0.25), Inches(2.9), Inches(5.2), Inches(1.8),
                 body, size=14, color=DARK)

    add_rounded_rect(s, Inches(0.8), Inches(5.2), Inches(11.9), Inches(1.35), LIGHT_GRAY)
    add_text(s, Inches(1.1), Inches(5.2), Inches(11.3), Inches(1.35),
             "명단을 고친 뒤에는 관리자에게 알려 주세요.\n"
             "이름 · 근무지는 로그인 화면에도 쓰이는 값이라 서버 쪽도 함께 맞춰야 합니다.",
             size=13.5, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 4)

    # ===== 5. 비밀번호 초기화 =====
    s = prs.slides.add_slide(blank)
    header(prs, s, "4. 비밀번호 초기화")

    add_text(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5),
             "작업자가 비밀번호를 잊었을 때의 절차입니다.",
             size=15, color=DARK, align=PP_ALIGN.CENTER)

    flow = [
        ("①", "[초기화 요청]", "명단에서 해당 작업자\n행의 버튼을 누릅니다"),
        ("②", "[초기화 대기]", "이름 옆에 표시가 붙습니다\n아직 초기화 전입니다"),
        ("③", "관리자 처리", "표시가 사라지면\n초기화된 것입니다"),
        ("④", "사번 + 1234", "작업자에게 알려 주세요\n로그인 후 새 비밀번호 등록"),
    ]
    for i, (num, title, desc) in enumerate(flow):
        x = Inches(0.7 + i * 3.15)
        add_rounded_rect(s, x, Inches(2.0), Inches(2.95), Inches(2.9), LIGHT_RED)
        add_text(s, x, Inches(2.2), Inches(2.95), Inches(0.6),
                 num, size=30, bold=True, color=RED, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(2.95), Inches(2.95), Inches(0.5),
                 title, size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(3.55), Inches(2.95), Inches(1.2),
                 desc, size=12, color=GRAY, align=PP_ALIGN.CENTER)

    add_rounded_rect(s, Inches(0.8), Inches(5.25), Inches(11.9), Inches(1.3), LIGHT_GRAY)
    add_text(s, Inches(1.1), Inches(5.25), Inches(11.3), Inches(1.3),
             "요청만으로는 초기화되지 않습니다. 관리자가 처리해야 [초기화 대기] 표시가 사라집니다.\n"
             "비밀번호는 서버가 복원할 수 없는 형태로 보관하므로 초기화만 가능합니다.",
             size=13.5, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 5)

    # ===== 6. 이럴 땐 =====
    s = prs.slides.add_slide(blank)
    header(prs, s, "5. 이럴 땐")

    qa = [
        ("[서버에서 불러오기] 버튼이 안 보임", "서무 사번으로 다시 로그인하세요"),
        ("미처리 휴가증이 없다고 나옴", "이미 가져왔거나, 7일 이후 사용 예정 건입니다"),
        ("[프로그램 실행] 을 눌러도 자동화가 안 뜸", "setup.bat 을 다시 실행해 주세요"),
        ("자동화가 중간에 멈춤", "error_*.png 와 last_error.log 를 관리자에게"),
        ("작업자가 잔여가 안 맞는다고 함", "손으로 고치지 말고 관리자에게 알려 주세요"),
        ("실수로 잘못된 휴가증을 처리함", "관리자에게 요청하면 지울 수 있습니다"),
    ]
    y0 = 1.35
    for i, (q, a) in enumerate(qa):
        y = Inches(y0 + i * 0.92)
        add_rounded_rect(s, Inches(0.8), y, Inches(11.9), Inches(0.78), LIGHT_GRAY)
        add_text(s, Inches(1.1), y, Inches(5.3), Inches(0.78),
                 q, size=13, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(6.6), y, Inches(5.9), Inches(0.78),
                 "→ " + a, size=13, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 6)

    prs.save(str(OUT))
    print(f"PPT 생성 완료: {OUT}")
    print(f"크기: {OUT.stat().st_size:,} bytes")


if __name__ == "__main__":
    main()
